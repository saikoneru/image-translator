from fastapi import FastAPI, UploadFile, Form, HTTPException
from fastapi.responses import Response
import httpx
from PIL import Image
import io, base64, json, os, tempfile, subprocess, asyncio
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from image_translator.utils import merge_translations_smart, draw_paragraphs_polys
from sse_starlette.sse import EventSourceResponse
import numpy as np

app = FastAPI(title="Unified Document Translator")
REQUEST_TIMEOUT = 300

WORKER_URLS = {
    "ocr": os.getenv("OCR_WORKER_URL", "http://i13hpc69:8001/ocr"),
    "segment": os.getenv("SEGMENT_WORKER_URL", "http://i13hpc69:8002/segment"),
    "inpaint": os.getenv("INPAINT_WORKER_URL", "http://i13hpc69:8003/inpaint"),
    "translate": os.getenv("TRANSLATE_WORKER_URL", "http://i13hpc69:8004/translate"),
    "multimodal_translate": os.getenv("MULTIMODAL_TRANSLATE_WORKER_URL", "http://i13hpc69:8006/multimodal_translate"),
    "layout": os.getenv("LAYOUT_WORKER_URL", "http://i13hpc69:8005/detect_layout"),
}

# ---------------------- IMAGE TRANSLATION ---------------------- #

async def translate_image_bytes(
    img_bytes: bytes,
    src_lang: str,
    tgt_lang: str,
    client: httpx.AsyncClient
) -> bytes:
    """Translate image by OCR, layout analysis, translation, and inpainting."""
    
    original = Image.open(io.BytesIO(img_bytes))
    has_transparency = original.mode in ("RGBA", "LA") or (
        original.mode == "P" and "transparency" in original.info
    )

    # 1. OCR
    ocr = (await client.post(
        WORKER_URLS["ocr"],
        files={"file": ("image.png", img_bytes)},
        timeout=REQUEST_TIMEOUT
    )).json()["results"]

    if not ocr or all(not (r.get("text") or "").strip() for r in ocr):
        return img_bytes

    # 2. Layout detection
    layout = (await client.post(
        WORKER_URLS["layout"],
        files={
            "file": ("image.png", img_bytes),
            "ocr_results_json": (None, json.dumps(ocr))
        },
        timeout=REQUEST_TIMEOUT
    )).json()

    # 3. Segment and collect texts
    all_texts = []
    mappings = []

    for para_idx, ocr_para in enumerate(layout.get("paragraphs", [])):
        seg = (await client.post(
            WORKER_URLS["segment"],
            files={"file": ("image.png", img_bytes)},
            data={"ocr_results_json": json.dumps(ocr_para), "join_all": True}
        )).json()

        merged = seg["merged_results"]
        start = len(all_texts)
        all_texts.extend([x["merged_text"] for x in merged])
        end = len(all_texts)

        mappings.append({
            "ocr_para": ocr_para,
            "merged": merged,
            "start": start,
            "end": end
        })

    # 4. Batch translation
    translated = []
    if all_texts:
        resp = await client.post(
            WORKER_URLS["multimodal_translate"],
            files={"file": ("image.png", img_bytes)},
            data={
                "texts_json": json.dumps(all_texts),
                "src_lang": src_lang,
                "tgt_lang": tgt_lang
            },
            timeout=60.0
        )
        resp.raise_for_status()
        translated = resp.json()["translations"]

    # 5. Merge translations back to line structure
    merged_paragraphs = []
    for m in mappings:
        block = translated[m["start"]:m["end"]]
        
        # Update merged results with translations
        merged_with_translations = []
        for entry, translation in zip(m["merged"], block):
            updated_entry = entry.copy()
            updated_entry["merged_text"] = translation
            merged_with_translations.append(updated_entry)

        # Merge back into OCR line structure
        ocr_para_copy = [line.copy() for line in m["ocr_para"]]
        merged_result = merge_translations_smart(merged_with_translations, ocr_para_copy)
        merged_paragraphs.append(merged_result)

    # 6. Inpaint original text regions
    inpaint = (await client.post(
        WORKER_URLS["inpaint"],
        files={"file": ("image.png", img_bytes)},
        data={"boxes_json": json.dumps([r["box"] for r in ocr])}
    )).json()

    inpainted = Image.open(io.BytesIO(base64.b64decode(inpaint["inpainted_image_base64"])))
    if has_transparency:
        inpainted = inpainted.convert("RGBA")

    # 7. Draw translated text
    base = original.convert("RGBA" if has_transparency else "RGB")
    result = draw_paragraphs_polys(inpainted.copy(), merged_paragraphs, base)

    buf = io.BytesIO()
    result.save(buf, format="PNG")
    return buf.getvalue()


# ---------------------- PPTX TRANSLATION HELPERS ---------------------- #

async def call_translate_api(texts, src, tgt, client):
    """Call translation API for text array."""
    if not texts:
        return []
    resp = await client.post(
        WORKER_URLS["translate"],
        data={"texts_json": json.dumps(texts), "src_lang": src, "tgt_lang": tgt}
    )
    data = resp.json()
    return data.get("translations") or data.get("results") or data.get("translated_texts") or []

def merge_translated_paragraph(paragraph, translated):
    """Replace paragraph text with translation."""
    full_translation = "\n".join(translated) if isinstance(translated, list) else translated
    paragraph.text = full_translation

async def process_text_frame(text_frame, src, tgt, client):
    """Translate all paragraphs in a text frame."""
    for p in text_frame.paragraphs:
        lines = p.text.split("\n")
        if any(lines):
            t = await call_translate_api(lines, src, tgt, client)
            merge_translated_paragraph(p, t)

async def process_table(table, src, tgt, client):
    """Translate all cells in a table."""
    for row in table.rows:
        for cell in row.cells:
            if hasattr(cell, "text_frame"):
                await process_text_frame(cell.text_frame, src, tgt, client)

async def translate_image_shape(shape, src, tgt, client):
    """Translate image embedded in shape."""
    try:
        orig = shape.image.blob
        new = await translate_image_bytes(orig, src, tgt, client)
        stream = io.BytesIO(new)
        image_part, rId = shape.part.get_or_add_image_part(stream)
        shape._element.blipFill.blip.rEmbed = rId
    except Exception:
        pass

async def process_shape(shape, src, tgt, client, translate_images=True):
    """Recursively process shape and translate content."""
    if getattr(shape, "has_text_frame", False):
        await process_text_frame(shape.text_frame, src, tgt, client)
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        await process_table(shape.table, src, tgt, client)
    elif getattr(shape, "image", None) and translate_images:
        await translate_image_shape(shape, src, tgt, client)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for shp in shape.shapes:
            await process_shape(shp, src, tgt, client, translate_images)

async def convert_pptx_to_pdf(pptx_bytes: bytes) -> bytes:
    """Convert PPTX to PDF via LibreOffice."""
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pptx")
        with open(input_path, "wb") as f:
            f.write(pptx_bytes)

        output_dir = os.path.join(tmpdir, "output")
        os.makedirs(output_dir, exist_ok=True)

        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pdf",
             "--outdir", output_dir, input_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120,
            check=False
        )

        if result.returncode != 0:
            raise Exception(f"LibreOffice conversion failed: {result.stderr.decode()}")

        pdf_path = os.path.join(output_dir, "input.pdf")
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF not generated at {pdf_path}")

        with open(pdf_path, "rb") as f:
            return f.read()

async def translate_pptx_bytes(
    pptx_bytes: bytes,
    src: str,
    tgt: str,
    client: httpx.AsyncClient,
    translate_master=True,
    translate_images=True
) -> bytes:
    """Translate PPTX file content."""
    prs = Presentation(io.BytesIO(pptx_bytes))

    # Translate masters and layouts
    if translate_master:
        masters = {id(slide.slide_layout.slide_master): slide.slide_layout.slide_master
                   for slide in prs.slides}.values()
        for m in masters:
            for shape in m.shapes:
                await process_shape(shape, src, tgt, client, translate_images)
            for layout in m.slide_layouts:
                for shape in layout.shapes:
                    await process_shape(shape, src, tgt, client, translate_images)

    # Translate slides
    for slide in prs.slides:
        for shape in slide.shapes:
            await process_shape(shape, src, tgt, client, translate_images)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------- ENDPOINTS ---------------------- #

@app.post("/translate/image")
async def translate_image(file: UploadFile, src_lang: str = Form(...), tgt_lang: str = Form(...)):
    """Translate an image file."""
    img = await file.read()
    async with httpx.AsyncClient() as c:
        out = await translate_image_bytes(img, src_lang, tgt_lang, c)
    return Response(content=out, media_type="image/png")

@app.post("/translate/pptx")
async def translate_pptx(
    file: UploadFile,
    src_lang: str = Form(...),
    tgt_lang: str = Form(...),
    output_format: str = Form("pdf"),
    translate_master: bool = Form(True),
    translate_images: bool = Form(True)
):
    """Translate a PPTX file and return as PDF or PPTX."""
    pptx = await file.read()
    async with httpx.AsyncClient() as client:
        translated = await translate_pptx_bytes(
            pptx, src_lang, tgt_lang, client, translate_master, translate_images
        )

    if output_format.lower() == "pdf":
        pdf = await convert_pptx_to_pdf(translated)
        return Response(
            content=pdf,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={file.filename.replace('.pptx','_translated.pdf')}"}
        )

    return Response(
        content=translated,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.post("/translate/auto")
async def translate_auto(
    file: UploadFile,
    src_lang: str = Form(...),
    tgt_lang: str = Form(...),
    output_format: str = Form("auto"),
    translate_master: bool = Form(True),
    translate_images: bool = Form(True)
):
    """Auto-detect file type and translate."""
    name = file.filename.lower()

    if name.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
        return await translate_image(file, src_lang, tgt_lang)

    if name.endswith(".pptx"):
        fmt = "pdf" if output_format == "auto" else output_format
        return await translate_pptx(file, src_lang, tgt_lang, fmt, translate_master, translate_images)

    raise HTTPException(400, "Unsupported file type")

@app.post("/translate/pptx/sse")
async def translate_pptx_sse(
    file: UploadFile,
    src_lang: str = Form(...),
    tgt_lang: str = Form(...),
    translate_master: bool = Form(True),
    translate_images: bool = Form(True)
):
    """Translate PPTX with Server-Sent Events progress updates."""
    pptx_bytes = await file.read()

    async def event_stream():
        try:
            prs = Presentation(io.BytesIO(pptx_bytes))
            total = len(prs.slides)

            yield {
                "event": "status",
                "data": json.dumps({"status": "started", "total_slides": total})
            }
            await asyncio.sleep(0)

            async with httpx.AsyncClient() as client:
                # Translate masters
                if translate_master:
                    yield {"event": "status", "data": json.dumps({"status": "processing_masters"})}
                    await asyncio.sleep(0)

                    masters = {id(slide.slide_layout.slide_master): slide.slide_layout.slide_master
                               for slide in prs.slides}.values()
                    for master in masters:
                        for shape in master.shapes:
                            await process_shape(shape, src_lang, tgt_lang, client, translate_images)
                        for layout in master.slide_layouts:
                            for shape in layout.shapes:
                                await process_shape(shape, src_lang, tgt_lang, client, translate_images)

                # Translate slides with progress
                for idx, slide in enumerate(prs.slides):
                    slide_no = idx + 1

                    yield {
                        "event": "progress",
                        "data": json.dumps({"slide": slide_no, "total": total})
                    }
                    await asyncio.sleep(0)

                    for shape in slide.shapes:
                        await process_shape(shape, src_lang, tgt_lang, client, translate_images)

                    # Convert to PDF and send
                    buf = io.BytesIO()
                    prs.save(buf)
                    pdf_bytes = await convert_pptx_to_pdf(buf.getvalue())
                    pdf_b64 = base64.b64encode(pdf_bytes).decode("utf-8")

                    yield {
                        "event": "slide",
                        "data": json.dumps({
                            "slide": slide_no,
                            "total": total,
                            "file_base64": pdf_b64,
                            "format": "pdf"
                        }),
                    }
                    await asyncio.sleep(0)

                yield {
                    "event": "complete",
                    "data": json.dumps({"status": "completed", "total_slides": total})
                }

        except Exception as e:
            yield {"event": "error", "data": json.dumps({"error": str(e)})}

    return EventSourceResponse(
        event_stream(),
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        ping=10
    )

@app.get("/health")
async def health():
    return {"status": "healthy"}
