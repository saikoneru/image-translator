import io
import base64
import requests
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
from PIL import Image
import gradio as gr
import json

# --- Defaults ---
DEFAULT_TRANSLATE_URL = "http://i13hpc69:8004/translate"
DEFAULT_IMAGE_PIPELINE_URL = "http://127.0.0.1:8080/process"
MIN_FONT_SIZE = 8

# ---------- Translation API ----------
def call_translate_api(texts, src_lang, tgt_lang, translate_url):
    if not texts:
        return []
    payload = {"texts_json": json.dumps(texts), "src_lang": src_lang, "tgt_lang": tgt_lang}
    resp = requests.post(translate_url, data=payload, timeout=300)
    resp.raise_for_status()
    data = resp.json()
    translations = data.get("translations") or data.get("results") or data.get("translated_texts")
    if translations is None:
        if isinstance(data, dict) and len(data) == 1:
            val = list(data.values())[0]
            if isinstance(val, str):
                return [val]
        raise RuntimeError(f"Unexpected translate API response: {data}")
    return translations

# ---------- Paragraph merging ----------
def merge_translated_paragraph_preserve_runs(paragraph, translated_lines, force_single_line=False):
    original_lines = paragraph.text.split("\n")
    num_lines = len(original_lines)

    # Adjust translated lines to match original line count
    if len(translated_lines) < num_lines:
        merged_text = " ".join(translated_lines)
        translated_lines = []
        start = 0
        for line in original_lines:
            line_len = len(line)
            translated_lines.append(merged_text[start:start+line_len])
            start += line_len
        if start < len(merged_text):
            translated_lines[-1] += merged_text[start:]
    elif len(translated_lines) > num_lines:
        merged_text = " ".join(translated_lines)
        avg_len = len(merged_text) // num_lines
        translated_lines = [merged_text[i*avg_len:(i+1)*avg_len] for i in range(num_lines)]
        if len(merged_text) % num_lines:
            translated_lines[-1] += merged_text[num_lines*avg_len:]

    # Distribute translated lines to runs
    run_texts = [r.text or "" for r in paragraph.runs]
    run_assignments = {i: [] for i in range(len(paragraph.runs))}
    run_idx = 0
    run_pos = 0
    for orig_line, trans_line in zip(original_lines, translated_lines):
        remaining = len(orig_line)
        segments = []
        while remaining > 0 and run_idx < len(run_texts):
            cur_text = run_texts[run_idx]
            avail = len(cur_text) - run_pos
            if avail <= 0:
                run_idx += 1
                run_pos = 0
                continue
            take = min(avail, remaining)
            segments.append((run_idx, cur_text[run_pos:run_pos+take]))
            remaining -= take
            run_pos += take
            if run_pos >= len(cur_text):
                run_idx += 1
                run_pos = 0
        # Assign translated line proportionally
        total_len = sum(len(s[1]) for s in segments) or 1
        pos = 0
        for i, (r_idx, s_text) in enumerate(segments):
            if i == len(segments) - 1:
                seg = trans_line[pos:]
            else:
                seg_len = int(round(len(s_text) / total_len * len(trans_line)))
                seg = trans_line[pos:pos+seg_len]
            run_assignments[r_idx].append(seg)
            pos += len(seg)

    for r_idx, parts in run_assignments.items():
        paragraph.runs[r_idx].text = "\n".join(parts) if parts else ""

# ---------- Text frame helpers ----------
def get_text_frame_width(text_frame):
    parent = getattr(text_frame, "_parent", None)
    if parent is None:
        return None
    if hasattr(parent, "width") and parent.width is not None:
        return parent.width.pt
    return None

def adjust_text_frame(text_frame, original_font_sizes):
    parent = getattr(text_frame, "_parent", None)
    if parent is None:
        return
    cur_width = get_text_frame_width(text_frame)
    if cur_width is None:
        return
    text_frame.auto_size = MSO_AUTO_SIZE.NONE

    # Shrink font if overflow
    for _ in range(20):
        overflow = False
        for paragraph in text_frame.paragraphs:
            total_text_len = sum(len(r.text) for r in paragraph.runs)
            avg_font_size = sum(r.font.size.pt for r in paragraph.runs if r.font.size) / max(len(paragraph.runs), 1)
            if total_text_len * avg_font_size > cur_width:
                overflow = True
                for run in paragraph.runs:
                    if run.font.size:
                        run.font.size = Pt(max(run.font.size.pt * 0.9, MIN_FONT_SIZE))
        if not overflow:
            break

    # Expand box width if space available
    if hasattr(parent, "width") and parent.width is not None:
        text_width_est = sum(len(run.text) * run.font.size.pt for paragraph in text_frame.paragraphs for run in paragraph.runs if run.font.size)
        if text_width_est > parent.width.pt:
            slide_width_pt = None
            if hasattr(text_frame._parent, "_parent") and hasattr(text_frame._parent._parent, "width") and text_frame._parent._parent.width is not None:
                slide_width_pt = text_frame._parent._parent.width.pt
            elif hasattr(parent, "_parent") and hasattr(parent._parent, "width") and parent._parent.width is not None:
                slide_width_pt = parent._parent.width.pt
            if slide_width_pt is None:
                slide_width_pt = parent.width.pt * 2
            max_width = max(1, slide_width_pt - getattr(parent, "left", Pt(0)).pt)
            new_width = min(text_width_est, max_width)
            parent.width = Pt(new_width)

# ---------- Process text and tables ----------
def process_text_frame(text_frame, src_lang, tgt_lang, translate_url, shrink_and_expand=True):
    original_font_sizes = [run.font.size for paragraph in text_frame.paragraphs for run in paragraph.runs]
    for paragraph in text_frame.paragraphs:
        lines = paragraph.text.split("\n")
        if not any(lines):
            continue
        is_single_line = len(lines) == 1
        try:
            translations = call_translate_api(lines, src_lang, tgt_lang, translate_url)
        except Exception as e:
            print("Translation API failed:", e)
            continue
        try:
            merge_translated_paragraph_preserve_runs(paragraph, translations, force_single_line=is_single_line)
        except Exception as e:
            paragraph.text = " ".join(translations) if is_single_line else "\n".join(translations)
    if shrink_and_expand:
        adjust_text_frame(text_frame, original_font_sizes)

def process_table(table, src_lang, tgt_lang, translate_url, shrink_and_expand=True):
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            if getattr(cell, "text_frame", None):
                process_text_frame(cell.text_frame, src_lang, tgt_lang, translate_url, shrink_and_expand)

# ---------- Image translation ----------
def replace_image_blob(shape, new_image_bytes):
    try:
        image_stream = io.BytesIO(new_image_bytes) if isinstance(new_image_bytes, bytes) else new_image_bytes
        image_part, rId = shape.part.get_or_add_image_part(image_stream)
        blip = shape._element.blipFill.blip
        blip.rEmbed = rId
    except Exception as e:
        print("replace_image_blob failed:", e)

def translate_image_shape(shape, src_lang, tgt_lang, image_pipeline_url):
    try:
        if not getattr(shape, "image", None):
            return
        orig_bytes = shape.image.blob
        img = Image.open(io.BytesIO(orig_bytes))
        has_transparency = img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info)
        original_format = img.format or 'PNG'
        if has_transparency and img.mode not in ('RGBA', 'LA'):
            img = img.convert('RGBA')
        img_buffer = io.BytesIO()
        img.save(img_buffer, format='PNG' if has_transparency else original_format)
        img_buffer.seek(0)
        send_bytes = img_buffer.getvalue()
        filename = "image.png" if has_transparency else f"image.{original_format.lower()}"
        mime_type = "image/png" if has_transparency else f"image/{original_format.lower()}"
        files = {"file": (filename, io.BytesIO(send_bytes), mime_type)}
        data = {"src_lang": src_lang, "tgt_lang": tgt_lang, "preserve_transparency": "true"}
        resp = requests.post(image_pipeline_url, files=files, data=data, timeout=600)
        resp.raise_for_status()
        img_b64 = resp.json().get("image_base64")
        if not img_b64:
            return
        new_bytes = base64.b64decode(img_b64)
        replace_image_blob(shape, new_bytes)
    except Exception as e:
        print("translate_image_shape failed:", e)

# ---------- Shape traversal ----------
def process_shape(shape, src_lang, tgt_lang, translate_url, image_pipeline_url, shrink_and_expand=True):
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        process_text_frame(shape.text_frame, src_lang, tgt_lang, translate_url, shrink_and_expand)
    elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TABLE and getattr(shape, "table", None):
        process_table(shape.table, src_lang, tgt_lang, translate_url, shrink_and_expand)
    elif getattr(shape, "image", None) is not None:
        translate_image_shape(shape, src_lang, tgt_lang, image_pipeline_url)
    elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for shp in shape.shapes:
            process_shape(shp, src_lang, tgt_lang, translate_url, image_pipeline_url, shrink_and_expand)

# ---------- Master/Layout ----------
def process_slide_master(slide_master, src_lang, tgt_lang, translate_url, image_pipeline_url):
    for shape in list(slide_master.shapes):
        process_shape(shape, src_lang, tgt_lang, translate_url, image_pipeline_url, shrink_and_expand=False)

def process_slide_layout(slide_layout, src_lang, tgt_lang, translate_url, image_pipeline_url):
    for shape in list(slide_layout.shapes):
        process_shape(shape, src_lang, tgt_lang, translate_url, image_pipeline_url, shrink_and_expand=False)

# ---------- PPTX processing ----------
def translate_pptx_bytes(input_bytes, src_lang, tgt_lang,
                         translate_master=True,
                         translate_url=DEFAULT_TRANSLATE_URL,
                         image_pipeline_url=DEFAULT_IMAGE_PIPELINE_URL):
    prs = Presentation(io.BytesIO(input_bytes))
    if translate_master:
        used_masters = []
        seen_master_ids = set()
        for slide in prs.slides:
            if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'slide_master'):
                master = slide.slide_layout.slide_master
                master_id = id(master)
                if master_id not in seen_master_ids:
                    seen_master_ids.add(master_id)
                    used_masters.append(master)
        for master in used_masters:
            process_slide_master(master, src_lang, tgt_lang, translate_url, image_pipeline_url)
            for layout in master.slide_layouts:
                process_slide_layout(layout, src_lang, tgt_lang, translate_url, image_pipeline_url)
    for slide in prs.slides:
        for shape in list(slide.shapes):
            process_shape(shape, src_lang, tgt_lang, translate_url, image_pipeline_url, shrink_and_expand=True)
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()

# ---------- Gradio ----------
def gradio_translate(pptx_file, src_lang, tgt_lang, translate_master, translate_url, image_pipeline_url):
    if pptx_file is None:
        return None
    with open(pptx_file.name, "rb") as f:
        in_bytes = f.read()
    out_bytes = translate_pptx_bytes(in_bytes, src_lang, tgt_lang, translate_master, translate_url, image_pipeline_url)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp.write(out_bytes)
    tmp.flush()
    tmp.close()
    return tmp.name

# ---------- Launch Gradio ----------
if __name__ == "__main__":
    langs = ["English", "Chinese", "Japanese", "French", "German"]
    with gr.Blocks() as demo:
        gr.Markdown("# PPTX Translator — Full Hierarchy & Run Preservation")
        with gr.Row():
            inp = gr.File(label="Upload PPTX", file_types=[".pptx"])
            out = gr.File(label="Download Translated PPTX")
        with gr.Row():
            src = gr.Dropdown(langs, value="English", label="Source language")
            tgt = gr.Dropdown(langs, value="German", label="Target language")
        translate_master_checkbox = gr.Checkbox(label="Translate slide masters and layouts", value=True)
        translate_url_input = gr.Textbox(DEFAULT_TRANSLATE_URL, label="Translate service URL")
        image_url_input = gr.Textbox(DEFAULT_IMAGE_PIPELINE_URL, label="Image pipeline URL (/process)")
        btn = gr.Button("Translate PPTX")
        btn.click(fn=gradio_translate,
                  inputs=[inp, src, tgt, translate_master_checkbox, translate_url_input, image_url_input],
                  outputs=out)
    demo.launch(server_name="0.0.0.0", server_port=7872, debug=False)
